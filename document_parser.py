# document_parser.py - Lit PDF, Word, PowerPoint

import PyPDF2
from docx import Document
from pptx import Presentation
import os


def lire_pdf(chemin: str) -> str:
    """Extrait tout le texte d'un PDF"""
    texte = ""
    try:
        with open(chemin, 'rb') as fichier:
            lecteur = PyPDF2.PdfReader(fichier)
            for i, page in enumerate(lecteur.pages):
                texte += f"\n--- Page {i+1} ---\n"
                texte += page.extract_text() + "\n"
        return texte
    except Exception as e:
        return f"❌ Erreur PDF : {e}"


def lire_word(chemin: str) -> str:
    """Extrait tout le texte d'un fichier Word"""
    try:
        doc = Document(chemin)
        texte = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        return texte
    except Exception as e:
        return f"❌ Erreur Word : {e}"


def lire_powerpoint(chemin: str) -> str:
    """Extrait tout le texte d'un PowerPoint"""
    try:
        prs = Presentation(chemin)
        texte = ""
        for i, slide in enumerate(prs.slides):
            texte += f"\n--- Slide {i+1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texte += shape.text + "\n"
        return texte
    except Exception as e:
        return f"❌ Erreur PPT : {e}"


def lire_document(chemin: str) -> str:
    """Lit n'importe quel document automatiquement"""
    if not os.path.exists(chemin):
        return f"❌ Fichier introuvable : {chemin}"
    
    extension = os.path.splitext(chemin)[1].lower()
    
    if extension == ".pdf":
        return lire_pdf(chemin)
    elif extension in [".docx", ".doc"]:
        return lire_word(chemin)
    elif extension in [".pptx", ".ppt"]:
        return lire_powerpoint(chemin)
    elif extension == ".txt":
        with open(chemin, 'r', encoding='utf-8') as f:
            return f.read()
    else:
        return f"❌ Format non supporté : {extension}"


# TEST
if __name__ == "__main__":
    print("📄 Module de lecture de documents prêt !")
    print("Formats supportés : PDF, Word (.docx), PowerPoint (.pptx), TXT")